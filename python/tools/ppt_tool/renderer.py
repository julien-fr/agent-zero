
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from .models import PresentationModel, TitleSlide, BulletSlide, QuoteSlide, ImageSlide, ChartSlide
from .theme import get_theme, ThemeSettings
import os

class PresentationRenderer:
    def __init__(self, data: PresentationModel):
        self.data = data
        self.theme = get_theme(data.theme)
        if data.template_path and os.path.exists(data.template_path):
            self.prs = Presentation(data.template_path)
        else:
            self.prs = Presentation()
            # Set slide width/height if needed, defaulting to 16:9
            self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _get_layout_safe(self, index):
        """Return the layout at index or the last available layout if index is out of bounds."""
        if index < len(self.prs.slide_layouts):
            return self.prs.slide_layouts[index]
        return self.prs.slide_layouts[-1]

    def _apply_background(self, slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        r, g, b = self.theme.colors.background
        fill.fore_color.rgb = RGBColor(r, g, b)

    def _apply_text_style(self, run, color_rgb, size=None, bold=False):
        run.font.name = self.theme.font_name
        r, g, b = color_rgb
        run.font.color.rgb = RGBColor(r, g, b)
        if size:
            run.font.size = size
        if bold:
            run.font.bold = True

    def render(self, output_path: str):
        # 1. Create Main Title Slide
        self._create_main_title_slide()

        # 2. Iterate content slides
        for slide_data in self.data.slides:
            if isinstance(slide_data, TitleSlide):
                self._create_section_header(slide_data)
            elif isinstance(slide_data, BulletSlide):
                self._create_bullet_slide(slide_data)
            elif isinstance(slide_data, QuoteSlide):
                self._create_quote_slide(slide_data)
            elif isinstance(slide_data, ImageSlide):
                self._create_image_slide(slide_data)
            elif isinstance(slide_data, ChartSlide):
                self._create_chart_slide(slide_data)

        self.prs.save(output_path)

    def _create_main_title_slide(self):
        slide_layout = self._get_layout_safe(0) # Title Slide
        slide = self.prs.slides.add_slide(slide_layout)
        self._apply_background(slide)

        # Safe Title Access
        try:
            title = slide.shapes.title
        except Exception:
            title = None

        if not title:
            # Fallback: Create title box
            title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))

        title.text = self.data.title
        for p in title.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                self._apply_text_style(run, self.theme.colors.title, Pt(54), True)

        if self.data.author:
            # Safe Subtitle Access
            if len(slide.placeholders) > 1:
                try:
                    subtitle = slide.placeholders[1]
                except Exception:
                    subtitle = slide.shapes.add_textbox(Inches(3), Inches(4.5), Inches(7), Inches(1))
            else:
                # Fallback: Create subtitle box
                subtitle = slide.shapes.add_textbox(Inches(3), Inches(4.5), Inches(7), Inches(1))
            title = slide.shapes.title
        else:
            # Fallback: Create title box
            title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))

        title.text = self.data.title
        for p in title.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                self._apply_text_style(run, self.theme.colors.title, Pt(54), True)

        if self.data.author:
            # Safe Subtitle Access
            if len(slide.placeholders) > 1:
                subtitle = slide.placeholders[1]
            else:
                # Fallback: Create subtitle box
                subtitle = slide.shapes.add_textbox(Inches(3), Inches(4.5), Inches(7), Inches(1))

            subtitle.text = self.data.author
            for p in subtitle.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    self._apply_text_style(run, self.theme.colors.text, Pt(24))

    def _create_section_header(self, data: TitleSlide):
        slide_layout = self._get_layout_safe(2) # Section Header
        slide = self.prs.slides.add_slide(slide_layout)
        self._apply_background(slide)

        # Safe Title
        if slide.shapes.title:
            title = slide.shapes.title
        else:
            title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))

        title.text = data.content
        for p in title.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                self._apply_text_style(run, self.theme.colors.title, Pt(44), True)

        if data.subcontent:
            # Safe Subtitle/Text Access
            if len(slide.placeholders) > 1:
                text = slide.placeholders[1]
            else:
                text = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9), Inches(1))

            text.text = data.subcontent
            for p in text.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    self._apply_text_style(run, self.theme.colors.text, Pt(24))

    def _create_bullet_slide(self, data: BulletSlide):
        slide_layout = self._get_layout_safe(1) # Title and Content
        slide = self.prs.slides.add_slide(slide_layout)
        self._apply_background(slide)

        # Safe Title
        if slide.shapes.title:
            title = slide.shapes.title
        else:
            title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))

        title.text = data.title
        for p in title.text_frame.paragraphs:
            for run in p.runs:
                self._apply_text_style(run, self.theme.colors.title, Pt(40), True)

        # Safe Body
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            if not body.has_text_frame:
                 # If placeholder exists but isn't text (e.g. picture), create new box
                 body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(5))
        else:
            body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(5))

        tf = body.text_frame
        tf.clear()  # Clear existing empty paragraphs
        tf.word_wrap = True

        for item in data.bullets:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            for run in p.runs:
                self._apply_text_style(run, self.theme.colors.text, Pt(24))

    def _create_quote_slide(self, data: QuoteSlide):
        slide_layout = self._get_layout_safe(6) # Blank or fallback
        slide = self.prs.slides.add_slide(slide_layout)
        self._apply_background(slide)

        # Centered text box for quote
        left = Inches(1)
        top = Inches(2)
        width = Inches(11.33)
        height = Inches(3)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.add_paragraph()
        p.text = f'"{data.text}"'
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            self._apply_text_style(run, self.theme.colors.accent, Pt(36), True)

        if data.author:
            p_auth = tf.add_paragraph()
            p_auth.text = f"- {data.author}"
            p_auth.alignment = PP_ALIGN.RIGHT
            for run in p_auth.runs:
                self._apply_text_style(run, self.theme.colors.text, Pt(20))

    def _create_image_slide(self, data: ImageSlide):
        slide_layout = self._get_layout_safe(6) # Blank or fallback
        slide = self.prs.slides.add_slide(slide_layout)
        self._apply_background(slide)

        if data.title:
            # Add title manually since it is blank layout
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            tf = title_box.text_frame
            p = tf.add_paragraph()
            p.text = data.title
            for run in p.runs:
                self._apply_text_style(run, self.theme.colors.title, Pt(40), True)

        # Add Image
        if os.path.exists(data.path):
            # Center image roughly
            top = Inches(2) if data.title else Inches(1)
            height = Inches(4.5)
            # We let width be auto-scaled to preserve aspect ratio, or we can calculate it.
            # Simple approach: add picture
            slide.shapes.add_picture(data.path, Inches(2), top, height=height)
        else:
            # Placeholder if image missing
            txBox = slide.shapes.add_textbox(Inches(4), Inches(3), Inches(5), Inches(1))
            txBox.text = f"[Image missing: {data.path}]"

        if data.caption:
            caption_box = slide.shapes.add_textbox(Inches(2), Inches(6.6), Inches(9), Inches(0.5))
            tf = caption_box.text_frame
            p = tf.add_paragraph()
            p.text = data.caption
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                self._apply_text_style(run, self.theme.colors.text, Pt(14), False)

    def _create_chart_slide(self, data: ChartSlide):
        slide_layout = self._get_layout_safe(5) # Title Only
        slide = self.prs.slides.add_slide(slide_layout)
        self._apply_background(slide)

        title = slide.shapes.title
        title.text = data.title
        for p in title.text_frame.paragraphs:
            for run in p.runs:
                self._apply_text_style(run, self.theme.colors.title, Pt(40), True)

        # Define chart data
        chart_data = CategoryChartData()
        chart_data.categories = data.categories
        for series in data.series:
            chart_data.add_series(series.name, series.values)

        # Determine chart type
        if data.chart_type == 'BAR_CLUSTERED':
            chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        elif data.chart_type == 'LINE':
            chart_type = XL_CHART_TYPE.LINE
        elif data.chart_type == 'PIE':
            chart_type = XL_CHART_TYPE.PIE
        else:
            chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED

        # Add chart
        x, y, cx, cy = Inches(2), Inches(2), Inches(9.33), Inches(4.5)
        graphic_frame = slide.shapes.add_chart(
            chart_type, x, y, cx, cy, chart_data
        )
        chart = graphic_frame.chart
